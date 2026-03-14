import os
import shutil
import glob
import subprocess
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import docx2txt
import pandas as pd
from PIL import Image
import pytesseract
import io

SOURCE_DIR = "."
OUTPUT_DIR = "_parsed"

# Ensure output directory exists and is clean
if not os.path.exists(OUTPUT_DIR):
    os.makedirs(OUTPUT_DIR)

def ensure_out_dir(rel_path):
    out_dir = os.path.join(OUTPUT_DIR, rel_path)
    os.makedirs(out_dir, exist_ok=True)
    return out_dir

def handle_pdf(file_path, rel_dir, filename):
    out_dir = ensure_out_dir(rel_dir)
    md_out = os.path.join(out_dir, f"{filename}.md")
    
    md_content = [f"# {filename}\n"]
    ocr_pages = 0
    try:
        doc = fitz.open(file_path)
        for page_num in range(len(doc)):
            page = doc[page_num]
            md_content.append(f"## Page {page_num + 1}\n")
            
            page_text = page.get_text().strip()
            
            # Extract images from PDF
            image_list = page.get_images(full=True)
            extracted_images = []
            for img_index, img in enumerate(image_list):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                image_ext = base_image["ext"]
                if image_ext not in ["png", "jpeg", "jpg"]:
                    image_ext = "png"
                img_name = f"{filename}_page{page_num+1}_img{img_index+1}.{image_ext}"
                img_out_path = os.path.join(out_dir, img_name)
                with open(img_out_path, "wb") as f:
                    f.write(image_bytes)
                extracted_images.append((img_out_path, img_name))
            
            # OCR fallback if page has very little text but has images
            if len(page_text) < 50 and extracted_images:
                if ocr_pages < 5:
                    print(f"  OCR processing page {page_num + 1}...")
                    ocr_text = ""
                    for img_path, _ in extracted_images:
                        try:
                            img_obj = Image.open(img_path)
                            text = pytesseract.image_to_string(img_obj, lang='chi_sim+eng')
                            if text.strip():
                                ocr_text += text + "\n"
                        except Exception as e:
                            print(f"OCR failed for {img_path}: {e}")
                    if ocr_text.strip():
                        page_text += "\n\n*(OCR Extracted Text)*\n" + ocr_text
                    ocr_pages += 1
                else:
                    page_text += "\n\n*(Further OCR Skipped, Image included below)*\n"

            if page_text:
                md_content.append(page_text + "\n")
            else:
                md_content.append("*(No text found on this page)*\n")

            for _, img_name in extracted_images:
                md_content.append(f"![Image](./{img_name})\n")
                
            md_content.append("\n---\n")
                    
        with open(md_out, "w", encoding="utf-8") as f:
            f.write("\n".join(md_content))
            
    except Exception as e:
        print(f"Error processing PDF {file_path}: {e}")

def extract_text_from_shape(shape):
    text = ""
    if hasattr(shape, "text") and shape.text:
        text += shape.text + "\n"
    if shape.has_table:
        for row in shape.table.rows:
            row_text = []
            for cell in row.cells:
                row_text.append(cell.text.replace('\n', ' ').strip())
            text += " | ".join(row_text) + "\n"
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for s in shape.shapes:
            text += extract_text_from_shape(s)
    return text

def handle_pptx(file_path, rel_dir, filename):
    out_dir = ensure_out_dir(rel_dir)
    md_out = os.path.join(out_dir, f"{filename}.md")
    
    md_content = [f"# {filename}\n"]
    img_count = 1
    try:
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            md_content.append(f"## Slide {i + 1}\n")
            
            slide_extracted_images = []
            
            def process_shape(shape):
                nonlocal img_count
                md_content.append(extract_text_from_shape(shape))
                
                # Extract embedded images
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = shape.image
                    image_bytes = image.blob
                    image_ext = image.ext
                    img_name = f"{filename}_slide{i+1}_img{img_count}.{image_ext}"
                    img_out_path = os.path.join(out_dir, img_name)
                    with open(img_out_path, "wb") as f:
                        f.write(image_bytes)
                    slide_extracted_images.append(img_name)
                    img_count += 1
                elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    for s in shape.shapes:
                        process_shape(s)

            for shape in slide.shapes:
                process_shape(shape)
                                
            # Notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                md_content.append("\n**Notes:**\n" + slide.notes_slide.notes_text_frame.text + "\n")
            
            for img_name in slide_extracted_images:
                md_content.append(f"![Slide Image](./{img_name})\n")
                
            md_content.append("\n---\n")
                
        with open(md_out, "w", encoding="utf-8") as f:
            f.write("\n".join(md_content))
            
    except Exception as e:
        print(f"Error processing PPTX {file_path}: {e}")

def handle_docx(file_path, rel_dir, filename):
    out_dir = ensure_out_dir(rel_dir)
    md_out = os.path.join(out_dir, f"{filename}.md")
    
    # Docx2txt extracts text AND saves images if image_dir is provided
    img_dir = os.path.join(out_dir, f"{filename}_images")
    os.makedirs(img_dir, exist_ok=True)
    
    try:
        text = docx2txt.process(file_path, img_dir)
        md_content = f"# {filename}\n\n"
        md_content += text if text else ""
        
        md_content += "\n\n## Extracted Images\n"
        if os.path.exists(img_dir):
            for img_file in os.listdir(img_dir):
                md_content += f"\n![Image](./{filename}_images/{img_file})"
                
        with open(md_out, "w", encoding="utf-8") as f:
            f.write(md_content)
            
        # Clean up empty image dir
        if not os.listdir(img_dir):
            os.rmdir(img_dir)
    except Exception as e:
        print(f"Error processing DOCX {file_path}: {e}")

def handle_doc(file_path, rel_dir, filename):
    out_dir = ensure_out_dir(rel_dir)
    md_out = os.path.join(out_dir, f"{filename}.md")
    
    try:
        result = subprocess.run(["antiword", file_path], capture_output=True, text=True)
        text = result.stdout
        md_content = f"# {filename}\n\n" + (text if text else "*(No text extracted)*")
        with open(md_out, "w", encoding="utf-8") as f:
            f.write(md_content)
    except Exception as e:
        print(f"Error processing DOC {file_path}: {e}")

def handle_heic(file_path, rel_dir, filename):
    out_dir = ensure_out_dir(rel_dir)
    jpg_out = os.path.join(out_dir, f"{filename}.jpg")
    
    try:
        subprocess.run(["sips", "-s", "format", "jpeg", file_path, "--out", jpg_out], check=True, stderr=subprocess.DEVNULL, stdout=subprocess.DEVNULL)
    except Exception as e:
        print(f"Error processing HEIC {file_path}: {e}")

def handle_excel(file_path, rel_dir, filename):
    out_dir = ensure_out_dir(rel_dir)
    md_out = os.path.join(out_dir, f"{filename}.md")
    
    try:
        df_dict = pd.read_excel(file_path, sheet_name=None)
        with open(md_out, "w", encoding="utf-8") as f:
            f.write(f"# {filename}\n\n")
            for sheet_name, df in df_dict.items():
                f.write(f"## Sheet: {sheet_name}\n\n")
                f.write(df.to_markdown(index=False) or "*(Empty Sheet)*")
                f.write("\n\n")
    except Exception as e:
        print(f"Error processing Excel {file_path}: {e}")

def process_files():
    total = 0
    for root, dirs, files in os.walk(SOURCE_DIR):
        if "_parsed" in root or ".venv" in root or "scripts" in root:
            continue
            
        rel_dir = os.path.relpath(root, SOURCE_DIR)
        if rel_dir == ".":
            rel_dir = ""
            
        for file in files:
            if file.startswith("~") or file.startswith("."):
                continue
                
            file_path = os.path.join(root, file)
            filename, ext = os.path.splitext(file)
            ext = ext.lower()
            
            print(f"Processing: {file_path}")
            total += 1
            
            if ext == ".pdf":
                handle_pdf(file_path, rel_dir, filename)
            elif ext == ".pptx":
                handle_pptx(file_path, rel_dir, filename)
            elif ext == ".docx":
                handle_docx(file_path, rel_dir, filename)
            elif ext == ".doc":
                handle_doc(file_path, rel_dir, filename)
            elif ext == ".heic":
                handle_heic(file_path, rel_dir, filename)
            elif ext in [".xls", ".xlsx"]:
                handle_excel(file_path, rel_dir, filename)
    print(f"Finished processing {total} files!")

if __name__ == "__main__":
    process_files()
